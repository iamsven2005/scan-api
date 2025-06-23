from fastapi import FastAPI, UploadFile, File,Depends, Header, HTTPException, Query
from presidio_analyzer import AnalyzerEngine
from fastapi.middleware.cors import CORSMiddleware
from sqlalchemy import create_engine, Column, Integer, String, DateTime, Text, text
from sqlalchemy.ext.declarative import declarative_base
from pgvector.sqlalchemy import Vector
from sqlalchemy.orm import sessionmaker, Session
from datetime import datetime, timedelta
import json
import os
from docx import Document
from openpyxl import load_workbook
import io
import zipfile
from PIL import Image
import pytesseract
import fitz  # PyMuPDF
from pptx import Presentation
import tempfile
import subprocess
from FlagEmbedding import FlagModel
from uuid import uuid4
from pydantic import BaseModel
from typing import Optional, List
from sqlalchemy.dialects.postgresql import ARRAY  # Optional if you use ARRAY instead
import torch
import torch.nn.functional as F
import logging
logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)

device = "cuda" if torch.cuda.is_available() else "cpu"
embedding_model = FlagModel("BAAI/bge-base-en", use_fp16=torch.cuda.is_available())

app = FastAPI()
analyzer = AnalyzerEngine()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Replace "*" with ["http://localhost"] for tighter control
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Database setup
POSTGRES_USER = os.getenv("POSTGRES_USER", "postgres")
POSTGRES_PASSWORD = os.getenv("POSTGRES_PASSWORD", "postgres")
POSTGRES_DB = os.getenv("POSTGRES_DB", "sensitivescan")
POSTGRES_HOST = os.getenv("POSTGRES_HOST", "db")
POSTGRES_PORT = os.getenv("POSTGRES_PORT", "5434")
DATABASE_URL = f"postgresql://{POSTGRES_USER}:{POSTGRES_PASSWORD}@{POSTGRES_HOST}:{POSTGRES_PORT}/{POSTGRES_DB}"
UPLOAD_DIR = "./uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)
engine = create_engine(DATABASE_URL)
with engine.connect() as conn:
    conn.execute(text("CREATE EXTENSION IF NOT EXISTS vector"))
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

class ScanLog(Base):
    __tablename__ = "scan_logs"
    id = Column(Integer, primary_key=True, index=True)
    timestamp = Column(DateTime, default=datetime.utcnow)
    file_name = Column(String, nullable=False)
    findings = Column(Text, nullable=False)
    embedding = Column(Vector(768))  # Use 384 for all-MiniLM-L6-v2

class ApiToken(Base):
    __tablename__ = "api_tokens"
    id = Column(Integer, primary_key=True, index=True)
    token = Column(String, unique=True, nullable=False, index=True)
    description = Column(String, nullable=True)
    created_at = Column(DateTime, default=datetime.utcnow)
    expires_at = Column(DateTime, nullable=True)
    is_active = Column(Integer, default=1)  # 1 = active, 0 = inactive

Base.metadata.create_all(bind=engine)

def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


def verify_token(
    x_api_token: str = Header(...),
    db: Session = Depends(get_db)
):
    token_record = db.query(ApiToken).filter(
        ApiToken.token == x_api_token,
        ApiToken.is_active == 1
    ).first()

    if not token_record:
        raise HTTPException(status_code=401, detail="Invalid or inactive token")

    if token_record.expires_at and token_record.expires_at < datetime.utcnow():
        raise HTTPException(status_code=401, detail="Token has expired")

    return token_record  # Optional: return token info for tracking

# Preseed API Token on startup
def preseed_api_token():
    db = SessionLocal()
    try:
        existing = db.query(ApiToken).filter(ApiToken.token == "PRESEEDED-TOKEN-123").first()
        if not existing:
            token = ApiToken(
                token="PRESEEDED-TOKEN-123",
                description="Default preseeded token",
                expires_at=datetime.utcnow() + timedelta(days=3650),  # valid for 10 years
                is_active=1
            )
            db.add(token)
            db.commit()
            print("✅ Preseeded API token created.")
        else:
            print("ℹ️ Preseeded token already exists.")
    finally:
        db.close()

# Call it after tables are created
Base.metadata.create_all(bind=engine)
preseed_api_token()


def extract_text_from_archive(file_bytes: bytes, format: str) -> str:
    extracted_text = ""
    with tempfile.TemporaryDirectory() as tmp_dir:
        archive_path = os.path.join(tmp_dir, f"archive.{format}")
        with open(archive_path, "wb") as f:
            f.write(file_bytes)

        try:
            if format == "rar":
                subprocess.run(["unrar", "x", "-y", archive_path, tmp_dir], check=True)
            elif format == "7z":
                subprocess.run(["7z", "x", "-y", f"-o{tmp_dir}", archive_path], check=True)
        except subprocess.CalledProcessError as e:
            return f"[ERROR] Failed to extract {format.upper()}: {e}"

        # Walk extracted files
        for root, _, files in os.walk(tmp_dir):
            for fname in files:
                fpath = os.path.join(root, fname)
                try:
                    with open(fpath, "rb") as f:
                        content = f.read()
                        if fname.endswith(".txt"):
                            extracted_text += content.decode("utf-8", errors="ignore") + "\n"
                        elif fname.endswith(".docx"):
                            extracted_text += extract_text_from_docx(content) + "\n"
                        elif fname.endswith(".xlsx"):
                            extracted_text += extract_text_from_xlsx(content) + "\n"
                        elif fname.endswith(".pdf"):
                            extracted_text += extract_text_from_pdf(content) + "\n"
                        elif fname.endswith(".pptx"):
                            extracted_text += extract_text_from_pptx(content) + "\n"
                        elif fname.endswith(('.png', '.jpg', '.jpeg')):
                            extracted_text += extract_text_from_image(content) + "\n"
                except Exception as e:
                    extracted_text += f"[ERROR reading {fname}]: {e}\n"

    return extracted_text

def extract_text_from_pptx(file_bytes: bytes) -> str:
    text = ""
    prs = Presentation(io.BytesIO(file_bytes))
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape_text = shape.text.strip()
                if shape_text:
                    text += shape_text + "\n"
    return text

def extract_text_from_pdf(file_bytes: bytes) -> str:
    text = ""
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text
def extract_text_from_docx(file_bytes: bytes) -> str:
    document = Document(io.BytesIO(file_bytes))
    return "\n".join([para.text for para in document.paragraphs])

def extract_text_from_xlsx(file_bytes: bytes, max_rows: int = 100) -> str:
    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    text = ""
    row_count = 0

    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join([str(cell) for cell in row if cell is not None])
            if row_text:
                text += row_text + "\n"
                row_count += 1
                if row_count >= max_rows:
                    return text  # stop once limit is reached

    return text

def extract_text_from_zip(file_bytes: bytes) -> str:
    zip_stream = io.BytesIO(file_bytes)
    extracted_text = ""
    with zipfile.ZipFile(zip_stream) as zf:
        for file_info in zf.infolist():
            if file_info.filename.endswith(('.txt', '.docx', '.xlsx')):
                with zf.open(file_info) as extracted_file:
                    inner_bytes = extracted_file.read()
                    if file_info.filename.endswith('.txt'):
                        extracted_text += inner_bytes.decode("utf-8", errors="ignore") + "\n"
                    elif file_info.filename.endswith('.docx'):
                        extracted_text += extract_text_from_docx(inner_bytes) + "\n"
                    elif file_info.filename.endswith('.pdf'):
                        extracted_text += extract_text_from_pdf(inner_bytes) + "\n"
                    elif file_info.filename.endswith('.xlsx'):
                        extracted_text += extract_text_from_xlsx(inner_bytes) + "\n"
                    elif file_info.filename.endswith('.pptx'):
                        extracted_text += extract_text_from_pptx(inner_bytes) + "\n"


    return extracted_text

def extract_text_from_image(file_bytes: bytes) -> str:
    try:
        image = Image.open(io.BytesIO(file_bytes))
        print(f"OCR Image Info: format={image.format}, size={image.size}")
        text = pytesseract.image_to_string(image)
        if not text.strip():
            print("[WARNING] OCR result is empty.")
        return text
    except Exception as e:
        print(f"[ERROR] OCR Exception: {str(e)}")
        return f"OCR Error: {str(e)}"

class TokenCreateRequest(BaseModel):
    description: Optional[str] = None
    days_valid: Optional[int] = 30

class TokenResponse(BaseModel):
    id: int
    token: str
    description: Optional[str]
    created_at: datetime
    expires_at: Optional[datetime]
    is_active: int

@app.get("/tokens", response_model=List[TokenResponse])
def list_tokens(db: Session = Depends(get_db)):
    tokens = db.query(ApiToken).all()
    return tokens

@app.post("/tokens", response_model=TokenResponse)
def create_token(req: TokenCreateRequest, db: Session = Depends(get_db)):
    new_token_str = str(uuid4())
    logger.info(new_token_str)
    expires_at = datetime.utcnow() + timedelta(days=req.days_valid) if req.days_valid else None

    token = ApiToken(
        token=new_token_str,
        description=req.description,
        expires_at=expires_at,
        is_active=1
    )
    db.add(token)
    db.commit()
    db.refresh(token)
    logger.info(token)

    return token

@app.delete("/tokens/{token_id}", response_model=dict)
def delete_token(token_id: int, db: Session = Depends(get_db)):
    token = db.query(ApiToken).filter(ApiToken.id == token_id).first()
    if not token:
        raise HTTPException(status_code=404, detail="Token not found")

    db.delete(token)
    db.commit()
    return {"detail": "Token deleted"}

@app.patch("/tokens/{token_id}", response_model=TokenResponse)
def toggle_token_status(token_id: int, is_active: int, db: Session = Depends(get_db)):
    token = db.query(ApiToken).filter(ApiToken.id == token_id).first()
    if not token:
        raise HTTPException(status_code=404, detail="Token not found")

    token.is_active = is_active
    db.commit()
    db.refresh(token)
    return token

@app.post("/scan")
async def scan_file(
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    token_record: ApiToken = Depends(verify_token)  # Protects with DB-based token
):
    content = await file.read()
    file_path = os.path.join(UPLOAD_DIR, file.filename)
    with open(file_path, "wb") as f:
        f.write(content)
    ext = file.filename.lower()

    if ext.endswith(".docx"):
        text = extract_text_from_docx(content)
    elif ext.endswith(".xlsx"):
        text = extract_text_from_xlsx(content)
    elif ext.endswith(".zip"):
        text = extract_text_from_zip(content)
    elif ext.endswith(".pdf"):
        text = extract_text_from_pdf(content)
    elif ext.endswith(".pptx"):
        text = extract_text_from_pptx(content)
    elif ext.endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff', '.webp')):
        text = extract_text_from_image(content)
    elif ext.endswith(".rar"):
        text = extract_text_from_archive(content, format="rar")
    elif ext.endswith(".7z"):
        text = extract_text_from_archive(content, format="7z")
    else:
        text = content.decode("utf-8", errors="ignore")

    print("Text", text)
    print("______________")

    # Embed the text
    prompt = "Represent this sentence for semantic search:"
    embedded = embedding_model.encode([f"{prompt} {text}"])
    embedded_tensor = torch.tensor(embedded).to(device)
    normalized_tensor = F.normalize(embedded_tensor, p=2, dim=1)
    embedding_vector = normalized_tensor[0].tolist()

    results = analyzer.analyze(text=text, language="en")
    findings = []
    for r in results[:10]:  # Limit to 10 findings
        finding_dict = r.to_dict()
        finding_dict["text"] = text[r.start:r.end]
        findings.append(finding_dict)

    db = SessionLocal()
    try:
        log_entry = ScanLog(
            file_name=file.filename,
            findings=json.dumps(findings),
            embedding=embedding_vector  # Save vector here
        )
        db.add(log_entry)
        db.commit()
    finally:
        db.close()

    return {"file": file.filename, "sensitive": bool(findings),
 "findings": findings}
@app.post("/scan/sensitive")
async def check_sensitive_only(
    file: UploadFile = File(...),
    token_record: ApiToken = Depends(verify_token)
):
    content = await file.read()
    file_path = os.path.join(UPLOAD_DIR, file.filename)
    with open(file_path, "wb") as f:
        f.write(content)
    ext = file.filename.lower()

    if ext.endswith(".docx"):
        text = extract_text_from_docx(content)
    elif ext.endswith(".xlsx"):
        text = extract_text_from_xlsx(content)
    elif ext.endswith(".zip"):
        text = extract_text_from_zip(content)
    elif ext.endswith(".pdf"):
        text = extract_text_from_pdf(content)
    elif ext.endswith(".pptx"):
        text = extract_text_from_pptx(content)
    elif ext.endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff', '.webp')):
        text = extract_text_from_image(content)
    elif ext.endswith(".rar"):
        text = extract_text_from_archive(content, format="rar")
    elif ext.endswith(".7z"):
        text = extract_text_from_archive(content, format="7z")
    else:
        text = content.decode("utf-8", errors="ignore")

    results = analyzer.analyze(text=text, language="en")
    is_sensitive = bool(results)

    return {"sensitive": is_sensitive}

@app.get("/scan")
def get_all_scans(
    db: Session = Depends(get_db),
):
    logs = db.query(ScanLog).order_by(ScanLog.timestamp.desc()).all()
    return [
        {
            "id": log.id,
            "timestamp": log.timestamp.isoformat(),
            "file_name": log.file_name,
            "findings": json.loads(log.findings),
            "sensitive": bool(json.loads(log.findings)),
            "embedding": list(log.embedding) if log.embedding else None  # <- fix here
        }
        for log in logs
    ]

@app.get("/search")
def semantic_search(
    query: str = Query(..., description="Search query to match against scanned documents"),
    db: Session = Depends(get_db),
    top_k: int = 5
):
    if not query.strip():
        raise HTTPException(status_code=400, detail="Query must not be empty")

    # Generate embedding for the query
    prompt = "Represent this sentence for semantic search:"
    embedded = embedding_model.encode([f"{prompt} {query}"])
    embedded_tensor = torch.tensor(embedded).to(device)
    normalized_query = F.normalize(embedded_tensor, p=2, dim=1)[0]

    # Retrieve all logs with embeddings
    logs = db.query(ScanLog).filter(ScanLog.embedding.isnot(None)).all()

    # Compute cosine similarities
    similarities = []
    for log in logs:
        doc_embedding = torch.tensor(log.embedding).to(device)
        doc_embedding = F.normalize(doc_embedding, p=2, dim=0)
        score = torch.dot(normalized_query, doc_embedding).item()
        similarities.append((score, log.file_name, log.id, log.timestamp.isoformat()))

    # Sort by score descending and return top_k
    similarities.sort(reverse=True, key=lambda x: x[0])
    return [
        {
            "score": round(score, 4),
            "file_name": filename,
            "id": id_,
            "timestamp": timestamp
        }
        for score, filename, id_, timestamp in similarities[:top_k]
    ]